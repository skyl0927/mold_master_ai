from pptx import Presentation
import re

pptx_path = r'I:\AI TEAM PJT\AI PROCESS MASTER\Feasibility Report Template.pptx'
prs = Presentation(pptx_path)

# 플레이스홀더 패턴
placeholder_pattern = r'<([^>]+)>'

with open('template_analysis.txt', 'w', encoding='utf-8') as f:
    f.write("=" * 80 + "\n")
    f.write("Feasibility Report Template 분석\n")
    f.write("=" * 80 + "\n\n")
    f.write(f"슬라이드 수: {len(prs.slides)}\n")
    f.write(f"슬라이드 크기: {prs.slide_width / 914400:.2f}\" x {prs.slide_height / 914400:.2f}\"\n\n")

    placeholders = set()

    for i, slide in enumerate(prs.slides):
        f.write(f"\n{'=' * 80}\n")
        f.write(f"슬라이드 {i + 1}\n")
        f.write(f"{'=' * 80}\n\n")

        for shape in slide.shapes:
            # 텍스트 Shape
            if hasattr(shape, 'text') and shape.text:
                text = shape.text
                found_placeholders = re.findall(placeholder_pattern, text)
                if found_placeholders:
                    placeholders.update(found_placeholders)
                    f.write(f"[{shape.name}] {text[:100]}\n")

            # 테이블 Shape
            if hasattr(shape, 'has_table') and shape.has_table:
                table = shape.table
                f.write(f"\n[테이블: {shape.name}] {len(table.rows)}행 x {len(table.columns)}열\n")

                for row_idx, row in enumerate(table.rows):
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                            found = re.findall(placeholder_pattern, cell_text)
                            if found:
                                placeholders.update(found)

                    if row_texts:
                        f.write(f"  행{row_idx}: {' | '.join(row_texts)}\n")
                f.write("\n")

    f.write("\n" + "=" * 80 + "\n")
    f.write("발견된 플레이스홀더 (<>로 표시된 입력 필드)\n")
    f.write("=" * 80 + "\n")
    for placeholder in sorted(placeholders):
        f.write(f"  - <{placeholder}>\n")

print("분석 완료: template_analysis.txt 파일을 확인하세요.")
